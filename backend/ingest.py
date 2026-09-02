"""
Turn a dropped file into something the Messages API can consume.

Two routes, chosen per file type:

  PDF   -> uploaded once via the Files API and referenced by file_id as a
           `document` block. This keeps the diagrams, tables, and figures
           intact, which matters more here than for a typical study app:
           an intact visual span means that is the channel that works.
  Other -> text extraction (docx / pptx / txt / md / csv), sent as a text block.
"""

from __future__ import annotations

import hashlib
import mimetypes
from dataclasses import dataclass, asdict
from pathlib import Path

MAX_TEXT_CHARS = 600_000  # ~150k tokens; refuse rather than silently truncate

PDF_EXT = {".pdf"}
DOCX_EXT = {".docx"}
PPTX_EXT = {".pptx"}
TEXT_EXT = {".txt", ".md", ".markdown", ".csv", ".rtf", ".text"}
IMAGE_EXT = {".png", ".jpg", ".jpeg", ".gif", ".webp"}

SUPPORTED = PDF_EXT | DOCX_EXT | PPTX_EXT | TEXT_EXT | IMAGE_EXT


class UnsupportedFile(Exception):
    pass


class FileTooLarge(Exception):
    pass


@dataclass
class Source:
    """One ingested course file."""
    name: str
    kind: str            # "pdf" | "image" | "text"
    sha: str
    file_id: str | None  # Files API id, for pdf/image
    text: str | None     # extracted text, for text-route files
    pages: int | None
    chars: int

    def to_dict(self) -> dict:
        d = asdict(self)
        # Don't ship the whole document body to the browser.
        d["text"] = None if self.text is None else self.text[:400]
        return d

    def as_content_blocks(self) -> list[dict]:
        """Content blocks for a user message."""
        if self.kind == "pdf":
            return [{
                "type": "document",
                "source": {"type": "file", "file_id": self.file_id},
                "title": self.name,
            }]
        if self.kind == "image":
            return [{
                "type": "image",
                "source": {"type": "file", "file_id": self.file_id},
            }]
        return [{
            "type": "text",
            "text": f"<course_file name=\"{self.name}\">\n{self.text}\n</course_file>",
        }]


# ---------------------------------------------------------------------------
# Text extraction
# ---------------------------------------------------------------------------

def _extract_docx(path: Path) -> str:
    import docx  # python-docx

    doc = docx.Document(str(path))
    out: list[str] = []
    for para in doc.paragraphs:
        if para.text.strip():
            style = (para.style.name or "").lower()
            prefix = "## " if "heading" in style else ""
            out.append(prefix + para.text.strip())
    for i, table in enumerate(doc.tables, 1):
        out.append(f"\n[table {i}]")
        for row in table.rows:
            cells = [c.text.strip().replace("\n", " ") for c in row.cells]
            out.append(" | ".join(cells))
    return "\n".join(out)


def _extract_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    out: list[str] = []
    for n, slide in enumerate(prs.slides, 1):
        out.append(f"\n--- Slide {n} ---")
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                out.append(shape.text_frame.text.strip())
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    out.append(" | ".join(c.text.strip() for c in row.cells))
        notes = getattr(slide, "notes_slide", None)
        if notes is not None and notes.notes_text_frame.text.strip():
            out.append("[speaker notes] " + notes.notes_text_frame.text.strip())
    return "\n".join(out)


def _extract_plain(path: Path) -> str:
    for enc in ("utf-8", "utf-16", "cp1252", "latin-1"):
        try:
            return path.read_text(encoding=enc)
        except (UnicodeDecodeError, UnicodeError):
            continue
    return path.read_bytes().decode("utf-8", errors="replace")


def _pdf_pages(path: Path) -> int | None:
    try:
        from pypdf import PdfReader

        return len(PdfReader(str(path)).pages)
    except Exception:
        return None


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------

def ingest(path: Path, client) -> Source:
    """Ingest one file. `client` is an anthropic.Anthropic instance."""
    ext = path.suffix.lower()
    if ext not in SUPPORTED:
        raise UnsupportedFile(
            f"{path.name}: {ext or 'no extension'} is not supported. "
            f"Use: {', '.join(sorted(SUPPORTED))}"
        )

    sha = hashlib.sha256(path.read_bytes()).hexdigest()[:16]

    if ext in PDF_EXT or ext in IMAGE_EXT:
        mime = mimetypes.guess_type(path.name)[0] or (
            "application/pdf" if ext in PDF_EXT else "application/octet-stream"
        )
        with path.open("rb") as fh:
            uploaded = client.beta.files.upload(file=(path.name, fh, mime))
        return Source(
            name=path.name,
            kind="pdf" if ext in PDF_EXT else "image",
            sha=sha,
            file_id=uploaded.id,
            text=None,
            pages=_pdf_pages(path) if ext in PDF_EXT else None,
            chars=0,
        )

    if ext in DOCX_EXT:
        text = _extract_docx(path)
    elif ext in PPTX_EXT:
        text = _extract_pptx(path)
    else:
        text = _extract_plain(path)

    text = text.strip()
    if not text:
        raise UnsupportedFile(f"{path.name}: no readable text found in the file.")
    if len(text) > MAX_TEXT_CHARS:
        # Deliberately not truncating - silently dropping half a lecture would
        # produce a quiz that looks complete and isn't.
        raise FileTooLarge(
            f"{path.name} is {len(text):,} characters, over the "
            f"{MAX_TEXT_CHARS:,} limit. Split it into two files and drop both."
        )

    return Source(
        name=path.name,
        kind="text",
        sha=sha,
        file_id=None,
        text=text,
        pages=None,
        chars=len(text),
    )
